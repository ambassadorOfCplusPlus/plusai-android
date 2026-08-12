"""Генерация офисных документов для агента Plus AI (python-pptx / python-docx).

Высокоуровневые функции: агент передаёт структуру (JSON), получаем файл.
Возвращают строку-статус (путь + что создано) или '[ошибка] ...'.
"""
import json
import os


def build_pptx(spec_json: str, out_path: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Pt
        spec = json.loads(spec_json)
        prs = Presentation()

        # Титульный слайд
        title = spec.get("title", "Презентация")
        subtitle = spec.get("subtitle", "")
        s0 = prs.slides.add_slide(prs.slide_layouts[0])
        s0.shapes.title.text = title
        if len(s0.placeholders) > 1:
            s0.placeholders[1].text = subtitle

        # Контентные слайды: заголовок + маркированный список
        for sl in spec.get("slides", []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sl.get("title", "")
            body = slide.placeholders[1].text_frame
            body.clear()
            bullets = sl.get("bullets", [])
            for i, b in enumerate(bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = str(b)
                p.font.size = Pt(18)

        os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
        prs.save(out_path)
        n = len(spec.get("slides", [])) + 1
        return f"презентация создана: {out_path} ({n} слайд(ов))"
    except Exception as e:  # noqa: BLE001
        return f"[ошибка] {type(e).__name__}: {e}"


def build_docx(spec_json: str, out_path: str) -> str:
    try:
        from docx import Document
        spec = json.loads(spec_json)
        doc = Document()

        if spec.get("title"):
            doc.add_heading(spec["title"], level=0)

        for sec in spec.get("sections", []):
            if sec.get("heading"):
                doc.add_heading(sec["heading"], level=1)
            text = sec.get("text", "")
            for para in str(text).split("\n"):
                if para.strip():
                    doc.add_paragraph(para)
            for b in sec.get("bullets", []):
                doc.add_paragraph(str(b), style="List Bullet")

        # Плоский вариант: просто параграфы без секций
        for para in spec.get("paragraphs", []):
            doc.add_paragraph(str(para))

        os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
        doc.save(out_path)
        return f"документ Word создан: {out_path}"
    except Exception as e:  # noqa: BLE001
        return f"[ошибка] {type(e).__name__}: {e}"
